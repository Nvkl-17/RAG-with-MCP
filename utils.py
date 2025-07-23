from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

# Load embedder once
embedder = SentenceTransformer("all-MiniLM-L6-v2")
index = None
doc_texts = []

def normalize(vecs):
    return vecs / np.linalg.norm(vecs, axis=1, keepdims=True)

def ingest_documents(files):
    global index, doc_texts
    texts = []

    for file in files:
        name = file.name.lower()
        if name.endswith(".pdf"):
            import fitz
            doc = fitz.open(stream=file.read(), filetype="pdf")
            texts.extend([page.get_text() for page in doc])
        elif name.endswith(".docx"):
            from docx import Document
            doc = Document(file)
            texts.extend([para.text for para in doc.paragraphs])
        elif name.endswith(".txt") or name.endswith(".md"):
            texts.append(file.read().decode("utf-8"))
        elif name.endswith(".csv"):
            import pandas as pd
            df = pd.read_csv(file)
            texts.append(df.to_string())
        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(file)
            for slide in prs.slides:
                slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                texts.append(slide_text)

    doc_texts.extend(texts)  # append to global list for multiple batches

    embeddings = embedder.encode(texts, convert_to_numpy=True)
    embeddings = normalize(embeddings)

    if index is None:
        index_dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(index_dim)  # Use cosine similarity
    index.add(embeddings)
    return index

def retrieve_chunks(query, kb_index, k=3):
    if kb_index is None:
        return []

    query_emb = embedder.encode([query], convert_to_numpy=True)
    query_emb = normalize(query_emb)

    distances, indices = index.search(query_emb, k)
    return [doc_texts[i] for i in indices[0]]

from transformers import AutoTokenizer, AutoModelForSeq2SeqLM

model_name = "google/flan-t5-base"
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModelForSeq2SeqLM.from_pretrained(model_name)

def generate_response(prompt, context_chunks):
    context = "\n".join(context_chunks)
    full_prompt = f"Context:\n{context}\n\nQuestion: {prompt}"
    inputs = tokenizer(full_prompt, return_tensors="pt", truncation=True, max_length=512)
    outputs = model.generate(**inputs, max_new_tokens=150)
    return tokenizer.decode(outputs[0], skip_special_tokens=True)
